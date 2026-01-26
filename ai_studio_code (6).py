import os
import re
import json
import sys
import time
import asyncio
import logging
import aiosqlite
import html
from datetime import datetime

from groq import AsyncGroq
from aiogram import Bot, Dispatcher, types, F
from aiogram.filters import Command, CommandObject
from aiogram.types import (ReplyKeyboardMarkup, KeyboardButton, InlineKeyboardMarkup,
                           InlineKeyboardButton, FSInputFile, CallbackQuery)
from aiogram.fsm.state import State, StatesGroup
from aiogram.fsm.context import FSMContext
from aiogram.client.default import DefaultBotProperties
from aiogram.enums import ParseMode
from aiohttp import web

# PPTX va Fayllar uchun
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import pypdf
from docx import Document

# --- 1. KONFIGURATSIYA ---
logging.basicConfig(level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s")
logger = logging.getLogger(__name__)

API_TOKEN = os.getenv('BOT_TOKEN')
GROQ_API_KEY = os.getenv('GROQ_API_KEY')
ADMIN_ID = int(os.getenv('ADMIN_ID', 0))
CHANNEL_ID = "@abdujalils" 

if not API_TOKEN or not GROQ_API_KEY:
    logger.critical("BOT_TOKEN yoki GROQ_API_KEY muhit o'zgaruvchilarida topilmadi!")
    sys.exit(1)

client = AsyncGroq(api_key=GROQ_API_KEY)
bot = Bot(token=API_TOKEN, default=DefaultBotProperties(parse_mode=ParseMode.HTML))
dp = Dispatcher()
DB_PATH = 'slide_master.db'

# --- 2. MULTILINGUAL CONTENT ---
LANGS = {
    'uz': {
        'welcome': "✨ <b>Slide Master AI Bot</b>\n\nSlaydlar, testlar va shpargalkalar tayyorlovchi aqlli yordamchi!",
        'btns': ["💎 Tariflar", "📊 Kabinet", "🤝 Taklif qilish", "❓ Quiz Test", "🌐 Til"],
        'sub_err': "🔒 <b>Botdan foydalanish uchun kanalga a'zo bo'ling:</b>",
        'tarif': "💎 <b>TARIFLAR:</b>\n\n⚡ 1 Slayd: 990 so'm\n🔥 5 Slayd: 2,999 so'm\n👑 VIP (Cheksiz): 5,999 so'm\n\n💳 Karta: <code>9860230107924485</code>\n📸 To'lovdan so'ng chekni yuboring:",
        'wait': "🧠 <b>AI ishlamoqda, kuting...</b>",
        'done': "✅ <b>Tayyor!</b>",
        'no_bal': "⚠️ Balans yetarli emas!",
        'cancel': "❌ Bekor qilish",
        'gen_prompt': "📝 <b>Mavzu:</b> {topic}\n\nNechta slayd kerak?",
        'quiz_prompt': "📂 <b>Fayl yuboring (PDF/DOCX/TXT):</b>\nMen undan test tuzib beraman.",
        'payment_sent': "✅ Chek qabul qilindi, kuting.",
        'package_btns': ["1️⃣ 1 ta Slayd", "5️⃣ 5 ta Slayd", "👑 VIP Premium"],
        'referral': "🎁 Taklif havolangiz:\n"
    },
    'ru': {
        'welcome': "✨ <b>Slide Master AI Bot</b>\n\nУмный помощник для создания слайдов и тестов!",
        'btns': ["💎 Тарифы", "📊 Кабинет", "🤝 Пригласить", "❓ Quiz Test", "🌐 Til"],
        'sub_err': "🔒 <b>Подпишитесь на канал, чтобы продолжить:</b>",
        'tarif': "💎 <b>ТАРИФЫ:</b>\n\n⚡ 1 Слайд: 990 сум\n🔥 5 Слайдов: 2,999 сум\n👑 VIP (Безлимит): 5,999 сум\n\n💳 Карта: <code>9860230107924485</code>\n📸 Отправьте чек после оплаты:",
        'wait': "🧠 <b>AI работает...</b>",
        'done': "✅ <b>Готово!</b>",
        'no_bal': "⚠️ Недостаточно баланса!",
        'cancel': "❌ Отмена",
        'gen_prompt': "📝 <b>Тема:</b> {topic}\n\nСколько слайдов нужно?",
        'quiz_prompt': "📂 <b>Отправьте файл (PDF/DOCX/TXT):</b>\nЯ создам тест.",
        'payment_sent': "✅ Чек получен, ожидайте.",
        'package_btns': ["1️⃣ 1 Слайд", "5️⃣ 5 Слайдов", "👑 VIP Premium"],
        'referral': "🎁 Ваша ссылка:\n"
    },
    'en': {
        'welcome': "✨ <b>Slide Master AI Bot</b>\n\nAI helper for slides, quizzes, and study materials!",
        'btns': ["💎 Pricing", "📊 Profile", "🤝 Invite", "❓ Quiz Test", "🌐 Til"],
        'sub_err': "🔒 <b>Please subscribe to our channel to use the bot:</b>",
        'tarif': "💎 <b>PRICING:</b>\n\n⚡ 1 Slide: 990 UZS\n🔥 5 Slides: 2,999 UZS\n👑 VIP (Unlimited): 5,999 UZS\n\n💳 Card: <code>9860230107924485</code>\n📸 Send us the receipt photo:",
        'wait': "🧠 <b>AI is processing...</b>",
        'done': "✅ <b>Done!</b>",
        'no_bal': "⚠️ Insufficient balance!",
        'cancel': "❌ Cancel",
        'gen_prompt': "📝 <b>Topic:</b> {topic}\n\nHow many slides?",
        'quiz_prompt': "📂 <b>Send a file (PDF/DOCX/TXT):</b>\nI will generate a quiz based on it.",
        'payment_sent': "✅ Receipt received, please wait.",
        'package_btns': ["1️⃣ 1 Slide", "5️⃣ 5 Slides", "👑 VIP Premium"],
        'referral': "🎁 Your invite link:\n"
    }
}

def get_text(lang, key):
    return LANGS.get(lang, LANGS['uz']).get(key, LANGS['uz'][key])

# --- 3. DATABASE MANAGER ---
class Database:
    def __init__(self, path):
        self.path = path

    async def init(self):
        async with aiosqlite.connect(self.path) as db:
            await db.execute("""CREATE TABLE IF NOT EXISTS users (
                id BIGINT PRIMARY KEY, username TEXT, lang TEXT DEFAULT 'uz', 
                is_premium INTEGER DEFAULT 0, balance INTEGER DEFAULT 2, 
                invited_by BIGINT, created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP)""")
            await db.execute("""CREATE TABLE IF NOT EXISTS payments (
                id INTEGER PRIMARY KEY AUTOINCREMENT, user_id BIGINT, 
                amount INTEGER, package_type TEXT, screenshot_id TEXT, 
                status TEXT DEFAULT 'pending', created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP)""")
            await db.commit()

    async def get_user(self, uid):
        async with aiosqlite.connect(self.path) as db:
            db.row_factory = aiosqlite.Row
            c = await db.execute("SELECT * FROM users WHERE id = ?", (uid,))
            return await c.fetchone()

    async def add_user(self, uid, username, ref_id=None):
        async with aiosqlite.connect(self.path) as db:
            try:
                await db.execute("INSERT INTO users (id, username, invited_by) VALUES (?, ?, ?)", (uid, username, ref_id))
                if ref_id:
                    await db.execute("UPDATE users SET balance = balance + 1 WHERE id = ?", (ref_id,))
                await db.commit()
                return True
            except: return False

    async def update_balance(self, uid, amount):
        async with aiosqlite.connect(self.path) as db:
            await db.execute("UPDATE users SET balance = balance + ? WHERE id = ?", (amount, uid))
            await db.commit()

    async def set_premium(self, uid):
        async with aiosqlite.connect(self.path) as db:
            await db.execute("UPDATE users SET is_premium = 1 WHERE id = ?", (uid,))
            await db.commit()

    async def set_lang(self, uid, lang):
        async with aiosqlite.connect(self.path) as db:
            await db.execute("UPDATE users SET lang = ? WHERE id = ?", (lang, uid))
            await db.commit()

    async def add_payment(self, uid, amount, pkg, photo):
        async with aiosqlite.connect(self.path) as db:
            c = await db.execute("INSERT INTO payments (user_id, amount, package_type, screenshot_id) VALUES (?, ?, ?, ?)", (uid, amount, pkg, photo))
            await db.commit()
            return c.lastrowid

    async def approve_payment(self, pid):
        async with aiosqlite.connect(self.path) as db:
            db.row_factory = aiosqlite.Row
            p = await (await db.execute("SELECT * FROM payments WHERE id = ?", (pid,))).fetchone()
            if p and p['status'] == 'pending':
                await db.execute("UPDATE payments SET status = 'approved' WHERE id = ?", (pid,))
                if p['package_type'] == 'vip_premium' or p['amount'] >= 999:
                    await db.execute("UPDATE users SET is_premium = 1 WHERE id = ?", (p['user_id'],))
                else:
                    await db.execute("UPDATE users SET balance = balance + ? WHERE id = ?", (p['amount'], p['user_id']))
                await db.commit()
                return p['user_id']
            return None

db = Database(DB_PATH)

# --- 4. STATES VA FAYL LOGIKASI ---
class UserStates(StatesGroup):
    choosing_pkg = State()
    waiting_pay = State()
    waiting_quiz = State()

def create_pptx(topic, json_data, uid):
    try:
        data = json.loads(json_data)
        prs = Presentation()
        prs.slide_width, prs.slide_height = Inches(13.33), Inches(7.5)
        
        for s_data in data.get('slides', []):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            # Fon
            bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
            bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(10, 15, 25); bg.line.width = 0
            
            # Sarlavha
            tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(1))
            p = tb.text_frame.paragraphs[0]
            p.text = s_data.get('title', topic).upper()
            p.font.bold = True; p.font.size = Pt(32); p.font.color.rgb = RGBColor(0, 210, 255)

            # Kontent
            cb = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(8.5), Inches(5.2))
            cb.fill.solid(); cb.fill.fore_color.rgb = RGBColor(25, 30, 45); cb.line.width = 0
            
            tf = cb.text_frame; tf.word_wrap = True
            for point in s_data.get('content', []):
                cp = tf.add_paragraph()
                cp.text = f"• {point}"; cp.font.size = Pt(18); cp.font.color.rgb = RGBColor(255, 255, 255); cp.space_before = Pt(10)

            # Insight
            if s_data.get('insight'):
                ib = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.2), Inches(1.5), Inches(3.6), Inches(3))
                ib.fill.solid(); ib.fill.fore_color.rgb = RGBColor(35, 45, 65); ib.line.color.rgb = RGBColor(255, 190, 0)
                itf = ib.text_frame; itf.word_wrap = True
                ip = itf.paragraphs[0]; ip.text = "💡 INFO"; ip.font.bold = True; ip.font.color.rgb = RGBColor(255, 190, 0)
                ip2 = itf.add_paragraph(); ip2.text = s_data['insight']; ip2.font.size = Pt(14); ip2.font.color.rgb = RGBColor(240, 240, 240)

        path = f"slides/S_{uid}_{int(time.time())}.pptx"
        os.makedirs("slides", exist_ok=True)
        prs.save(path)
        return path
    except: return None

# --- 5. HANDLERLAR ---
async def check_sub(uid):
    try:
        m = await bot.get_chat_member(CHANNEL_ID, uid)
        return m.status in ['member', 'administrator', 'creator']
    except: return True

async def main_menu(msg, lang):
    b = get_text(lang, 'btns')
    kb = ReplyKeyboardMarkup(keyboard=[
        [KeyboardButton(text=b[0]), KeyboardButton(text=b[1])],
        [KeyboardButton(text=b[2]), KeyboardButton(text=b[3])],
        [KeyboardButton(text=b[4])]
    ], resize_keyboard=True)
    await msg.answer(get_text(lang, 'welcome'), reply_markup=kb)

@dp.message(Command("start"))
async def cmd_start(msg: types.Message, cmd: CommandObject, state: FSMContext):
    await state.clear()
    uid = msg.from_user.id
    ref_id = int(cmd.args) if cmd.args and cmd.args.isdigit() and int(cmd.args) != uid else None
    await db.add_user(uid, msg.from_user.username, ref_id)
    u = await db.get_user(uid)
    if not await check_sub(uid):
        kb = InlineKeyboardMarkup(inline_keyboard=[
            [InlineKeyboardButton(text="📢 Join Channel", url=f"https://t.me/{CHANNEL_ID[1:]}")],
            [InlineKeyboardButton(text="✅ Check", callback_data="check_sub")]
        ])
        return await msg.answer(get_text(u['lang'], 'sub_err'), reply_markup=kb)
    await main_menu(msg, u['lang'])

@dp.callback_query(F.data == "check_sub")
async def cb_sub(cb: CallbackQuery):
    if await check_sub(cb.from_user.id):
        await cb.message.delete()
        u = await db.get_user(cb.from_user.id)
        await main_menu(cb.message, u['lang'])
    else: await cb.answer("❌ Sub first!", show_alert=True)

# Tilni o'zgartirish
@dp.callback_query(F.data.startswith("setlang_"))
async def cb_lang(cb: CallbackQuery):
    new_lang = cb.data.split("_")[1]
    await db.set_lang(cb.from_user.id, new_lang)
    await cb.message.delete()
    await main_menu(cb.message, new_lang)

# Asosiy tugmalar
@dp.message(F.text)
async def handle_text(msg: types.Message, state: FSMContext):
    u = await db.get_user(msg.from_user.id)
    if not u: return
    l, txt = u['lang'], msg.text
    
    # Tugmalarni aniqlash
    if txt in [LANGS['uz']['btns'][0], LANGS['ru']['btns'][0], LANGS['en']['btns'][0]]: # Tariflar
        pkgs = get_text(l, 'package_btns')
        kb = ReplyKeyboardMarkup(keyboard=[[KeyboardButton(text=pkgs[0]), KeyboardButton(text=pkgs[1])], [KeyboardButton(text=pkgs[2])], [KeyboardButton(text=get_text(l, 'cancel'))]], resize_keyboard=True)
        await msg.answer(get_text(l, 'tarif'), reply_markup=kb)
        await state.set_state(UserStates.choosing_pkg)
    
    elif txt in [LANGS['uz']['btns'][1], LANGS['ru']['btns'][1], LANGS['en']['btns'][1]]: # Kabinet
        status = "👑 VIP" if u['is_premium'] else "🆓 Basic"
        await msg.answer(f"👤 ID: {u['id']}\n💰 Bal: {u['balance']}\n🏷 Stat: {status}")

    elif txt in [LANGS['uz']['btns'][2], LANGS['ru']['btns'][2], LANGS['en']['btns'][2]]: # Invite
        me = await bot.get_me()
        await msg.answer(f"{get_text(l, 'referral')}<code>https://t.me/{me.username}?start={u['id']}</code>")

    elif txt in [LANGS['uz']['btns'][3], LANGS['ru']['btns'][3], LANGS['en']['btns'][3]]: # Quiz
        await msg.answer(get_text(l, 'quiz_prompt'), reply_markup=ReplyKeyboardMarkup(keyboard=[[KeyboardButton(text=get_text(l, 'cancel'))]], resize_keyboard=True))
        await state.set_state(UserStates.waiting_quiz)

    elif txt in [LANGS['uz']['btns'][4], LANGS['ru']['btns'][4], LANGS['en']['btns'][4]]: # Til
        kb = InlineKeyboardMarkup(inline_keyboard=[
            [InlineKeyboardButton(text="O'zbekcha 🇺🇿", callback_data="setlang_uz")],
            [InlineKeyboardButton(text="Русский 🇷🇺", callback_data="setlang_ru")],
            [InlineKeyboardButton(text="English 🇺🇸", callback_data="setlang_en")]
        ])
        await msg.answer("Choose language:", reply_markup=kb)

    elif txt in [get_text('uz', 'cancel'), get_text('ru', 'cancel'), get_text('en', 'cancel')]:
        await state.clear(); await main_menu(msg, l)

    else: # Mavzu (Slayd yaratish)
        if not u['is_premium'] and u['balance'] <= 0: return await msg.answer(get_text(l, 'no_bal'))
        await state.update_data(topic=txt)
        kb = InlineKeyboardMarkup(inline_keyboard=[[InlineKeyboardButton(text="7", callback_data="gen:7"), InlineKeyboardButton(text="10", callback_data="gen:10"), InlineKeyboardButton(text="15", callback_data="gen:15")]])
        await msg.answer(get_text(l, 'gen_prompt').format(topic=html.escape(txt)), reply_markup=kb)

# Slayd Generatsiya
@dp.callback_query(F.data.startswith("gen:"))
async def cb_gen(cb: CallbackQuery, state: FSMContext):
    uid = cb.from_user.id
    u = await db.get_user(uid)
    count = cb.data.split(":")[1]
    data = await state.get_data()
    topic = data.get('topic')
    
    await cb.message.delete()
    wm = await cb.message.answer(get_text(u['lang'], 'wait'))
    
    try:
        sys_p = "You are a Presentation AI. Return ONLY JSON. Format: {'slides': [{'title': '...', 'content': ['...', '...'], 'insight': '...'}]}"
        prompt = f"Topic: {topic}. Slides: {count}. Language: {u['lang']}."
        
        res = await client.chat.completions.create(
            messages=[{"role":"system","content":sys_p}, {"role":"user","content":prompt}],
            model="llama-3.3-70b-versatile", response_format={"type": "json_object"}
        )
        path = await asyncio.to_thread(create_pptx, topic, res.choices[0].message.content, uid)
        
        if path:
            await bot.send_document(uid, FSInputFile(path), caption=get_text(u['lang'], 'done'))
            if not u['is_premium']: await db.update_balance(uid, -1)
            os.remove(path)
        else: raise Exception()
    except: await cb.message.answer("❌ AI Error")
    finally: await wm.delete(); await state.clear()

# Quiz Handler
@dp.message(UserStates.waiting_quiz, F.document)
async def quiz_doc(msg: types.Message, state: FSMContext):
    u = await db.get_user(msg.from_user.id)
    wm = await msg.answer(get_text(u['lang'], 'wait'))
    path = f"temp_{msg.from_user.id}"
    await bot.download(msg.document, destination=path)
    
    try:
        # Fayldan matn olish
        ext = msg.document.file_name.split('.')[-1].lower()
        if ext == 'pdf': txt = "\n".join([p.extract_text() for p in pypdf.PdfReader(path).pages])
        elif ext == 'docx': txt = "\n".join([p.text for p in Document(path).paragraphs])
        else: txt = open(path, 'r', encoding='utf-8').read()
        
        prompt = f"Based on this text, create 10 multiple choice questions in {u['lang']}:\n\n{txt[:10000]}"
        res = await client.chat.completions.create(messages=[{"role":"user","content":prompt}], model="llama-3.3-70b-versatile")
        await msg.answer(res.choices[0].message.content)
    except: await msg.answer("❌ Error reading file")
    finally:
        if os.path.exists(path): os.remove(path)
        await wm.delete(); await state.clear(); await main_menu(msg, u['lang'])

# To'lov state
@dp.message(UserStates.choosing_pkg)
async def pkg_sel(msg: types.Message, state: FSMContext):
    u = await db.get_user(msg.from_user.id)
    p = get_text(u['lang'], 'package_btns')
    amt = 1 if msg.text == p[0] else 5 if msg.text == p[1] else 999
    await state.update_data(amt=amt, pkg=msg.text)
    await msg.answer("📸 Send receipt photo:")
    await state.set_state(UserStates.waiting_pay)

@dp.message(UserStates.waiting_pay, F.photo)
async def pay_photo(msg: types.Message, state: FSMContext):
    d = await state.get_data()
    pid = await db.add_payment(msg.from_user.id, d['amt'], d['pkg'], msg.photo[-1].file_id)
    kb = InlineKeyboardMarkup(inline_keyboard=[[InlineKeyboardButton(text="✅ OK", callback_data=f"p_ok_{pid}"), InlineKeyboardButton(text="❌ NO", callback_data=f"p_no_{pid}")]])
    await bot.send_photo(ADMIN_ID, msg.photo[-1].file_id, caption=f"Pay ID: {pid}\nUser: {msg.from_user.id}", reply_markup=kb)
    u = await db.get_user(msg.from_user.id)
    await msg.answer(get_text(u['lang'], 'payment_sent'))
    await state.clear(); await main_menu(msg, u['lang'])

@dp.callback_query(F.data.startswith("p_"))
async def adm_p(cb: CallbackQuery):
    act, pid = cb.data.split("_")[1], int(cb.data.split("_")[2])
    if act == "ok":
        uid = await db.approve_payment(pid)
        if uid:
            try: await bot.send_message(uid, "✅ Payment Approved!")
            except: pass
            await cb.message.edit_caption(caption="✅ Approved")
    else: await cb.message.edit_caption(caption="❌ Rejected")

# --- 6. SERVER VA RUN ---
async def health(req): return web.Response(text="Bot OK")
async def start_srv():
    app = web.Application()
    app.router.add_get('/', health)
    r = web.AppRunner(app)
    await r.setup()
    await web.TCPSite(r, '0.0.0.0', int(os.getenv("PORT", 8080))).start()

async def main():
    await db.init()
    asyncio.create_task(start_srv())
    await dp.start_polling(bot)

if __name__ == "__main__":
    asyncio.run(main())