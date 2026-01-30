import html
import logging
import asyncio
import os
import re
import json
import sys
import time
import aiosqlite
from groq import AsyncGroq
from aiogram import Bot, Dispatcher, types, F
from aiogram.filters import Command, CommandObject
from aiogram.types import (ReplyKeyboardMarkup, KeyboardButton, InlineKeyboardMarkup,
                           InlineKeyboardButton, FSInputFile, CallbackQuery)
from aiogram.fsm.state import State, StatesGroup
from aiogram.fsm.context import FSMContext
from aiogram.client.default import DefaultBotProperties
from aiogram.enums import ParseMode, ContentType
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Fayllarni o'qish uchun kutubxonalar
import pypdf
from docx import Document

# --- 1. KONFIGURATSIYA VA LOGGING ---
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s - %(name)s - %(levelname)s - %(message)s",
    handlers=[
        logging.StreamHandler(),
        logging.FileHandler('bot.log', encoding='utf-8')
    ]
)
logger = logging.getLogger(__name__)

# Environment variable'lardan o'qish
API_TOKEN = os.getenv('BOT_TOKEN')
GROQ_API_KEY = os.getenv('GROQ_API_KEY')
ADMIN_ID = os.getenv('ADMIN_ID')
CHANNEL_ID = "@abdujalils" # Kanal usernameni o'zgartiring

# Xavfsizlik tekshiruvi
if not API_TOKEN:
    logger.critical("❌ BOT_TOKEN sozlanmagan!")
    sys.exit(1)

if not GROQ_API_KEY:
    logger.critical("❌ GROQ_API_KEY sozlanmagan!")
    sys.exit(1)

try:
    ADMIN_ID = int(ADMIN_ID) if ADMIN_ID else 0
except ValueError:
    logger.warning("⚠️ ADMIN_ID noto'g'ri formatda, admin funksiyalari ishlamaydi.")
    ADMIN_ID = 0

# Global obyektlar
client = AsyncGroq(api_key=GROQ_API_KEY)
bot = Bot(token=API_TOKEN, default=DefaultBotProperties(parse_mode=ParseMode.MARKDOWN))
dp = Dispatcher()
DB_PATH = 'slide_master.db'

# --- 2. HOLATLAR (STATES) ---
class UserStates(StatesGroup):
    waiting_package_choice = State() # Paket tanlashni kutish
    waiting_for_payment = State()    # To'lov chekini kutish
    waiting_for_quiz_file = State()  # Quiz uchun fayl kutish
    
class AdminStates(StatesGroup):
    waiting_for_broadcast = State()

# --- 3. MULTILINGUAL KONTENT ---
MARKETING_TEXT = """🔥 **TALABALAR UCHUN SHOK: BITTA ISSIQ NON NARXIGA CHEKSIZ PREZENTATSIYALAR!** 🥖🤖

Sessiya yaqin, "mustaqil ish"lar esa tog'dek yig'ilib ketdimi? Kechalari uxlamasdan slayd yasashdan charchadingizmi? 😫

Sizga aqlbovar qilmaydigan yangilik! **Slide Master AI Bot** — endi prezentatsiya qilish siz uchun dunyodagi eng oson ishga aylanadi.

🤯 **NARXGA QARANG (Haqiqiy talababop):**
Kodimizdagi narxlarni shunchaki "tekin" darajasiga tushirdik:

🔹 **1 dona slayd** — atigi **990 so'm** (Hatto avtobus yo'lkirasidan ham arzon! 🚌)
🔹 **5 dona slayd** — **2,999 so'm**
🔹 **VIP PREMIUM (CHEKSIZ)** — atigi **5,999 so'm!** 🔥🔥🔥
*(Ha, adashmadingiz! Bor-yo'g'i bitta ISSIQ NON narxiga umrbod cheksiz prezentatsiyalar yarating!)*

💡 **Slide Master AI nima qila oladi?**
✅ **60 soniyada tayyor fayl:** Shunchaki mavzuni yozing, Llama-3 AI sizga professional tuzilmani yaratib beradi.
✅ **3 xil tilda:** O'zbek, Rus va Ingliz tillarida mukammal mantiq.
✅ **Formatlash:** Tayyor .pptx (PowerPoint) faylini yuklab oling va ishingizni topshiring.

🎁 **TEKIN FOYDALANISH YO'LLARI:**
1️⃣ Dastlabki **2 ta slayd** — ro'yxatdan o'tganingiz uchun mutlaqo BEPUL!
2️⃣ **Do'stlarni taklif qiling:** Referal havolangiz orqali kelgan har bir do'stingiz uchun **+1 bepul slayd** oling!

⚠️ *Diqqat: 5,999 so'mlik VIP narxi faqat dastlabki foydalanuvchilar uchun! Narxlar tez orada ko'tarilishi mumkin.*

**Vaqtingizni tejang, prezentatsiyalarni AI ga topshiring! 👇**

🚀 **Havola:** """

LANGS = {
    'uz': {
        'welcome': "✨ **Slide Master AI Bot**\n\nProfessional taqdimotlar yaratuvchi sun'iy intellekt!\n\n👇 Quyidagi menyudan kerakli bo'limni tanlang:",
        'btns': ["💎 Tariflar", "📊 Kabinet", "🤝 Taklif qilish", "❓ Quiz Test", "🌐 Til / Language"],
        'sub_err': "🔒 **Botdan foydalanish cheklangan!**\n\nDavom etish uchun rasmiy kanalimizga obuna bo'ling:",
        'tarif': "💎 **TAQDIMOT NARXLARI:**\n\n⚡ **1 ta Slayd:** 990 so'm\n🔥 **5 ta Slayd:** 2,999 so'm\n👑 **VIP Premium (Cheksiz):** 5,999 so'm\n\n💳 **To'lov kartasi:** `9860230107924485`\n👤 **Karta egasi:** Abdujalil A.\n\n📸 *Paketni tanlang va keyin to'lov chekini yuboring:*",
        'choose_package': "🛒 **Paketni tanlang:**",
        'wait': "🧠 **AI ishlamoqda...**\n\nSlayd tuzilishi va dizayni generatsiya qilinmoqda. Bu jarayon mavzu murakkabligiga qarab 30-60 soniya vaqt oladi.",
        'done': "✅ **Taqdimot tayyor!**\n\nFaylni ochish uchun PowerPoint yoki WPS Office ishlating.",
        'no_bal': "⚠️ **Balans yetarli emas!**\n\nSizda bepul urinishlar tugadi. Hisobni to'ldiring yoki do'stlaringizni taklif qiling.",
        'cancel': "❌ Bekor qilish",
        'lang_name': "🇺🇿 O'zbekcha",
        'gen_prompt': "Mavzu: {topic}. Nechta slayd kerak?",
        'btn_check': "✅ Obunani tekshirish",
        'btn_join': "📢 Kanalga qo'shilish",
        'error': "⚠️ Xatolik yuz berdi. Iltimos qayta urinib ko'ring yoki keyinroq urining.",
        'payment_sent': "✅ Chek adminga yuborildi. Tez orada javob beriladi.\n\n📋 *To'lov tasdiqlangandan so'ng siz tanlagan paket aktivlashtiriladi.*",
        'admin_panel': "🛠 **Admin panel**\n\nTanlang:",
        'help_text': "📚 **QO'LLANMA**\n\n1️⃣ Botdan foydalanish uchun kanalga obuna bo'ling\n2️⃣ Mavzu yozing va slayd sonini tanlang\n3️⃣ AI siz uchun prezentatsiya yaratadi\n4️⃣ Faylni PowerPoint yoki WPS Office'da oching\n\n🤝 Taklif qilingan har bir do'st uchun +1 slayd bonus!",
        'package_btns': ["1️⃣ 1 ta Slayd", "5️⃣ 5 ta Slayd", "👑 VIP Premium"],
        'balance_added': "💰 **Balans to'ldirildi!**\n\nHisobingizga **{amount} ta slayd** qo'shildi!",
        'premium_activated': "👑 **Tabriklaymiz!**\nSiz VIP Premium (cheksiz) statusga o'tdingiz!\nEndi cheksiz slayd yaratishingiz mumkin!",
        'send_check_now': "📸 **Endi to'lov chekini rasm sifatida yuboring:**",
        'quiz_prompt': "📂 **Faylni yuboring!**\n\nMen faylni o'qib, undagi ma'lumotlardan test (quiz) tuzib beraman.\n\n📄 Formatlar: **PDF, DOCX, TXT**",
        'quiz_processing': "⏳ **Fayl o'qilmoqda va test tuzilmoqda...**",
        'quiz_error': "⚠️ Faylni o'qishda xatolik bo'ldi. Matnli PDF yoki Word fayl yuboring."
    },
    'ru': {
        'welcome': "✨ **Slide Master AI Bot**\n\nИскусственный интеллект для создания профессиональных презентаций!\n\n👇 Выберите раздел из меню ниже:",
        'btns': ["💎 Тарифы", "📊 Кабинет", "🤝 Пригласить", "❓ Quiz Test", "🌐 Til / Language"],
        'sub_err': "🔒 **Доступ ограничен!**\n\nПожалуйста, подпишитесь на наш канал для продолжения:",
        'tarif': "💎 **ТАРИФЫ НА ПРЕЗЕНТАЦИИ:**\n\n⚡ **1 Слайд:** 990 сум\n🔥 **5 Слайдов:** 2,999 сум\n👑 **VIP Premium (Безлимит):** 5,999 сум\n\n💳 **Карта для оплаты:** `9860230107924485`\n👤 **Владелец карты:** Abdujalil A.\n\n📸 *Выберите пакет, затем отправьте скриншот чека:*",
        'choose_package': "🛒 **Выберите пакет:**",
        'wait': "🧠 **AI работает...**\n\nГенерируем структуру и дизайн. Это может занять 30-60 секунд.",
        'done': "✅ **Презентация готова!**\n\nИспользуйте PowerPoint или WPS Office для открытия.",
        'no_bal': "⚠️ **Недостаточно баланса!**\n\nПополните счет или приглашайте друзей для бесплатных слайдов.",
        'cancel': "❌ Отмена",
        'lang_name': "🇷🇺 Русский",
        'gen_prompt': "Тема: {topic}. Сколько слайдов нужно?",
        'btn_check': "✅ Проверить подписку",
        'btn_join': "📢 Подписаться",
        'error': "⚠️ Произошла ошибка. Пожалуйста, попробуйте снова.",
        'payment_sent': "✅ Чек отправлен администратору. Скоро получите ответ.\n\n📋 *После подтверждения оплаты выбранный пакет будет активирован.*",
        'admin_panel': "🛠 **Админ панель**\n\nВыберите:",
        'help_text': "📚 **ИНСТРУКЦИЯ**\n\n1️⃣ Подпишитесь на канал для использования бота\n2️⃣ Напишите тему и выберите количество слайдов\n3️⃣ AI создаст презентацию\n4️⃣ Откройте файл в PowerPoint или WPS Office\n\n🤝 +1 слайд бонус за каждого приглашенного друга!",
        'package_btns': ["1️⃣ 1 Слайд", "5️⃣ 5 Слайдов", "👑 VIP Premium"],
        'balance_added': "💰 **Баланс пополнен!**\n\nНа ваш счет добавлено **{amount} слайдов**!",
        'premium_activated': "👑 **Поздравляем!**\nВы перешли на VIP Premium (безлимит) статус!\nТеперь вы можете создавать неограниченное количество слайдов!",
        'send_check_now': "📸 **Теперь отправьте фото чека об оплате:**",
        'quiz_prompt': "📂 **Отправьте файл!**\n\nЯ прочитаю файл и создам тест (квиз) на основе информации.\n\n📄 Форматы: **PDF, DOCX, TXT**",
        'quiz_processing': "⏳ **Читаю файл и создаю тест...**",
        'quiz_error': "⚠️ Ошибка при чтении файла. Отправьте текстовый PDF или Word."
    },
    'en': {
        'welcome': "✨ **Slide Master AI Bot**\n\nProfessional presentation generator powered by AI!\n\n👇 Choose a section from the menu below:",
        'btns': ["💎 Pricing", "📊 Profile", "🤝 Invite", "❓ Quiz Test", "🌐 Til / Language"],
        'sub_err': "🔒 **Access Restricted!**\n\nPlease subscribe to our channel to continue:",
        'tarif': "💎 **PRESENTATION PRICING:**\n\n⚡ **1 Slide:** 990 UZS\n🔥 **5 Slides:** 2,999 UZS\n👑 **VIP Premium (Unlimited):** 5,999 UZS\n\n💳 **Payment card:** `9860230107924485`\n👤 **Card owner:** Abdujalil A.\n\n📸 *Choose a package and then send the receipt screenshot:*",
        'choose_package': "🛒 **Choose package:**",
        'wait': "🧠 **AI is thinking...**\n\nGenerating structure and design. This may take 30-60 seconds.",
        'done': "✅ **Presentation ready!**\n\nUse PowerPoint or WPS Office to open.",
        'no_bal': "⚠️ **Insufficient balance!**\n\nTop up your account or invite friends to get free slides.",
        'cancel': "❌ Cancel",
        'lang_name': "🇬🇧 English",
        'gen_prompt': "Topic: {topic}. How many slides needed?",
        'btn_check': "✅ Check Subscription",
        'btn_join': "📢 Join Channel",
        'error': "⚠️ An error occurred. Please try again.",
        'payment_sent': "✅ Receipt sent to admin. You'll get response soon.\n\n📋 *After payment confirmation, your chosen package will be activated.*",
        'admin_panel': "🛠 **Admin Panel**\n\nSelect:",
        'help_text': "📚 **GUIDE**\n\n1️⃣ Subscribe to channel to use bot\n2️⃣ Write topic and select slide count\n3️⃣ AI will create presentation\n4️⃣ Open file in PowerPoint or WPS Office\n\n🤝 +1 slide bonus for each invited friend!",
        'package_btns': ["1️⃣ 1 Slide", "5️⃣ 5 Slides", "👑 VIP Premium"],
        'balance_added': "💰 **Balance topped up!**\n\n**{amount} slides** added to your account!",
        'premium_activated': "👑 **Congratulations!**\nYou have upgraded to VIP Premium (unlimited) status!\nNow you can create unlimited slides!",
        'send_check_now': "📸 **Now send the payment receipt as a photo:**",
        'quiz_prompt': "📂 **Send a file!**\n\nI will read the file and create a quiz based on the information.\n\n📄 Formats: **PDF, DOCX, TXT**",
        'quiz_processing': "⏳ **Reading file and generating quiz...**",
        'quiz_error': "⚠️ Error reading file. Please send a text-based PDF or Word file."
    }
}

def get_text(lang_code, key):
    """Xavfsiz matn olish funksiyasi"""
    return LANGS.get(lang_code, LANGS['uz']).get(key, LANGS['uz'].get(key, "Text not found"))

# --- 4. BAZA MANAGER ---
class Database:
    def __init__(self, db_path):
        self.db_path = db_path

    async def init(self):
        async with aiosqlite.connect(self.db_path) as db:
            await db.execute("""
                CREATE TABLE IF NOT EXISTS users (
                    id BIGINT PRIMARY KEY,
                    username TEXT,
                    first_name TEXT,
                    last_name TEXT,
                    lang TEXT DEFAULT 'uz',
                    is_premium INTEGER DEFAULT 0,
                    balance INTEGER DEFAULT 2,
                    invited_by BIGINT,
                    created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                    last_active TIMESTAMP DEFAULT CURRENT_TIMESTAMP
                )
            """)
            await db.execute("""
                CREATE TABLE IF NOT EXISTS referrals (
                    id INTEGER PRIMARY KEY AUTOINCREMENT,
                    referrer_id BIGINT,
                    referred_id BIGINT,
                    bonus_given INTEGER DEFAULT 0,
                    created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
                )
            """)
            await db.execute("""
                CREATE TABLE IF NOT EXISTS payments (
                    id INTEGER PRIMARY KEY AUTOINCREMENT,
                    user_id BIGINT,
                    amount INTEGER,
                    package_type TEXT,
                    screenshot_id TEXT,
                    status TEXT DEFAULT 'pending',
                    created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
                )
            """)
            await db.commit()

    async def get_user(self, user_id):
        async with aiosqlite.connect(self.db_path) as db:
            db.row_factory = aiosqlite.Row
            cursor = await db.execute("SELECT * FROM users WHERE id = ?", (user_id,))
            return await cursor.fetchone()

    async def add_user(self, user_id, username, first_name, last_name, referrer_id=None):
        async with aiosqlite.connect(self.db_path) as db:
            try:
                await db.execute("""
                    INSERT INTO users (id, username, first_name, last_name, invited_by, balance) 
                    VALUES (?, ?, ?, ?, ?, 2)
                """, (user_id, username, first_name, last_name, referrer_id))
                await db.commit()
                
                if referrer_id:
                    await db.execute("""
                        INSERT INTO referrals (referrer_id, referred_id) 
                        VALUES (?, ?)
                    """, (referrer_id, user_id))
                    await db.commit()
                return True
            except aiosqlite.IntegrityError:
                await db.execute("UPDATE users SET last_active = CURRENT_TIMESTAMP WHERE id = ?", (user_id,))
                await db.commit()
                return False

    async def update_balance(self, user_id, amount):
        async with aiosqlite.connect(self.db_path) as db:
            await db.execute("UPDATE users SET balance = balance + ? WHERE id = ?", (amount, user_id))
            await db.commit()

    async def set_premium(self, user_id):
        async with aiosqlite.connect(self.db_path) as db:
            await db.execute("UPDATE users SET is_premium = 1 WHERE id = ?", (user_id,))
            await db.commit()

    async def update_lang(self, user_id, lang):
        async with aiosqlite.connect(self.db_path) as db:
            await db.execute("UPDATE users SET lang = ? WHERE id = ?", (lang, user_id))
            await db.commit()

    async def get_referral_count(self, user_id):
        async with aiosqlite.connect(self.db_path) as db:
            cursor = await db.execute("SELECT COUNT(*) FROM referrals WHERE referrer_id = ?", (user_id,))
            res = await cursor.fetchone()
            return res[0] if res else 0

    async def get_all_users(self):
        async with aiosqlite.connect(self.db_path) as db:
            db.row_factory = aiosqlite.Row
            cursor = await db.execute("SELECT id FROM users")
            return await cursor.fetchall()

    async def get_stats(self):
        async with aiosqlite.connect(self.db_path) as db:
            cursor = await db.execute("SELECT COUNT(*) as total_users, SUM(balance) as total_slides FROM users")
            stats = await cursor.fetchone()
            cursor2 = await db.execute("SELECT COUNT(*) as premium_users FROM users WHERE is_premium = 1")
            premium_stats = await cursor2.fetchone()
            
            return {
                'total_users': stats[0] if stats else 0,
                'total_slides': stats[1] if stats else 0,
                'premium_users': premium_stats[0] if premium_stats else 0
            }

    async def add_payment(self, user_id, amount, package_type, screenshot_id):
        async with aiosqlite.connect(self.db_path) as db:
            cursor = await db.execute("""
                INSERT INTO payments (user_id, amount, package_type, screenshot_id)
                VALUES (?, ?, ?, ?)
            """, (user_id, amount, package_type, screenshot_id))
            await db.commit()
            return cursor.lastrowid

db = Database(DB_PATH)

# --- 5. FAYL O'QISH FUNKSIYALARI ---
def extract_text_from_file(file_path):
    ext = file_path.split('.')[-1].lower()
    text = ""
    try:
        if ext == 'pdf':
            reader = pypdf.PdfReader(file_path)
            for page in reader.pages:
                text += page.extract_text() + "\n"
        elif ext == 'docx':
            doc = Document(file_path)
            for para in doc.paragraphs:
                text += para.text + "\n"
        elif ext == 'txt':
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()
    except Exception as e:
        logger.error(f"Fayl o'qishda xato: {e}")
        return None
    
    # Juda uzun matn bo'lsa qisqartiramiz (AI limiti uchun)
    return text[:15000] if text else None

# --- 6. PPTX GENERATOR ---
def clean_json_string(text):
    text = text.strip()
    json_match = re.search(r'```(?:json)?\s*(\{.*?\})\s*```', text, re.DOTALL)
    if json_match:
        return json_match.group(1)
    
    start = text.find('{')
    end = text.rfind('}')
    if start != -1 and end != -1:
        return text[start:end+1]
    return text

def create_ultra_modern_pptx(topic, json_data, uid):
    try:
        cleaned_json = clean_json_string(json_data)
        data = json.loads(cleaned_json)

        prs = Presentation()
        prs.slide_width = Inches(13.333)  # 16:9 keng format
        prs.slide_height = Inches(7.5)

        # YANGI ZAMONAVIY RANG PALITRASI
        BG_COLOR = RGBColor(13, 17, 23)      # To'q kosmik ko'k
        ACCENT_NEON = RGBColor(0, 247, 255)  # Neon moviy
        ACCENT_CORAL = RGBColor(255, 95, 109) # Korall qizil
        ACCENT_LIME = RGBColor(202, 255, 112) # Yashil-limon
        TEXT_WHITE = RGBColor(245, 247, 250) # Sof oq
        TEXT_GRAY = RGBColor(170, 180, 200)  # Kulrang
        CARD_BG = RGBColor(22, 27, 34)       # Kartalar fon
        GRADIENT_START = RGBColor(30, 40, 60) # Gradient boshlanish
        GRADIENT_END = RGBColor(15, 20, 30)   # Gradient tugashi

        for idx, s_data in enumerate(data.get('slides', [])):
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            # GRADIENT FON
            bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
            bg.fill.gradient()
            bg.fill.gradient_stops[0].color.rgb = GRADIENT_START
            bg.fill.gradient_stops[1].color.rgb = GRADIENT_END
            bg.line.fill.background()

            # GEOMETRIK DEKORATSIYALAR
            # 1. Chap tomonda diagonal chiziq
            line1 = slide.shapes.add_shape(MSO_SHAPE.LINE, Inches(0), Inches(2), Inches(4), Inches(0))
            line1.line.color.rgb = ACCENT_NEON
            line1.line.width = Pt(2)
            line1.rotation = -15
            
            # 2. O'ng tomonda nuqtalar
            for i in range(5):
                dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, 
                    Inches(11 + i*0.3), Inches(0.5 + i*0.5), 
                    Inches(0.1), Inches(0.1))
                dot.fill.solid()
                dot.fill.fore_color.rgb = ACCENT_CORAL
                dot.line.fill.background()

            # SARDALHA - YANGI DIZAYN
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(1.2))
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = s_data.get('title', topic).upper()
            p.font.size = Pt(36)  # Kattaroq shrift
            p.font.bold = True
            p.font.color.rgb = TEXT_WHITE
            p.font.name = "Calibri"
            
            # Sarlavha ostidagi dekorativ chiziq
            title_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 
                Inches(0.5), Inches(1.3), Inches(3), Inches(0.08))
            title_line.fill.solid()
            title_line.fill.fore_color.rgb = ACCENT_NEON
            title_line.line.fill.background()

            # SUBTITLE (agar mavjud bo'lsa)
            sub = s_data.get('subtitle', '')
            if sub:
                sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(10), Inches(0.5))
                sp = sub_box.text_frame.paragraphs[0]
                sp.text = sub
                sp.font.size = Pt(16)
                sp.font.italic = True
                sp.font.color.rgb = ACCENT_LIME

            # ASOSIY KONTEYNER - YANGI DIZAYN
            main_container = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, 
                Inches(0.5), Inches(2.0), Inches(8.5), Inches(4.8)
            )
            main_container.fill.solid()
            main_container.fill.fore_color.rgb = CARD_BG
            main_container.line.color.rgb = ACCENT_NEON
            main_container.line.width = Pt(1.5)
            main_container.shadow.inherit = False
            main_container.shadow.blur_radius = Pt(10)
            main_container.shadow.offset_x = Pt(2)
            main_container.shadow.offset_y = Pt(2)
            main_container.shadow.color.rgb = RGBColor(0, 0, 0)

            # KONTEYNER ICHIDAGI KONTENT
            content_tf = main_container.text_frame
            content_tf.word_wrap = True
            content_tf.margin_left = Inches(0.25)
            content_tf.margin_top = Inches(0.25)
            content_tf.margin_right = Inches(0.25)
            content_tf.margin_bottom = Inches(0.25)
            
            points = s_data.get('content', [])
            if isinstance(points, list) and points and isinstance(points[0], str):
                temp = []
                for t in points: 
                    temp.append({'bold': 'Asosiy', 'text': t})
                points = temp

            for i, point in enumerate(points):
                if i == 0: 
                    p = content_tf.paragraphs[0]
                else: 
                    p = content_tf.add_paragraph()
                
                # ZAMONAVIY BULLET POINT
                run_bullet = p.add_run()
                run_bullet.text = "◆ "
                run_bullet.font.color.rgb = ACCENT_CORAL
                run_bullet.font.size = Pt(18)
                run_bullet.font.bold = True

                bold_txt = point.get('bold', '')
                if bold_txt:
                    run_bold = p.add_run()
                    run_bold.text = f"{bold_txt}: "
                    run_bold.font.bold = True
                    run_bold.font.color.rgb = ACCENT_LIME
                    run_bold.font.size = Pt(18)

                run_main = p.add_run()
                run_main.text = point.get('text', '')
                run_main.font.color.rgb = TEXT_WHITE
                run_main.font.size = Pt(18)
                p.space_after = Pt(16)
                p.line_spacing = 1.2

            # STATISTIKA BLOGI - YANGI DIZAYN
            stat_val = s_data.get('stat', '')
            if stat_val:
                # Statistikani o'rab turuvchi dekorativ ramka
                stat_decor = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE, 
                    Inches(9.2), Inches(2.0), Inches(3.5), Inches(2.0)
                )
                stat_decor.fill.solid()
                stat_decor.fill.fore_color.rgb = RGBColor(30, 35, 42)
                stat_decor.line.color.rgb = ACCENT_CORAL
                stat_decor.line.width = Pt(2)
                stat_decor.line.dash_style = 2  # Chiziqli
                
                # Statistik qiymat
                stat_box = slide.shapes.add_textbox(Inches(9.3), Inches(2.3), Inches(3.3), Inches(1.0))
                stat_tf = stat_box.text_frame
                stat_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                p_stat = stat_tf.paragraphs[0]
                p_stat.text = str(stat_val)
                p_stat.font.size = Pt(42)
                p_stat.font.bold = True
                p_stat.font.color.rgb = ACCENT_NEON
                p_stat.alignment = PP_ALIGN.CENTER
                
                # Statistik sarlavha
                stat_label = slide.shapes.add_textbox(Inches(9.3), Inches(3.5), Inches(3.3), Inches(0.5))
                label_p = stat_label.text_frame.paragraphs[0]
                label_p.text = "📊 ASOSIY KO'RSATKICH"
                label_p.font.size = Pt(11)
                label_p.font.color.rgb = TEXT_GRAY
                label_p.alignment = PP_ALIGN.CENTER

            # INSIGHT BLOGI - YANGI DIZAYN
            insight_val = s_data.get('insight', '')
            if insight_val:
                insight_bg = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE, 
                    Inches(9.2), Inches(4.2), Inches(3.5), Inches(2.6)
                )
                insight_bg.fill.solid()
                insight_bg.fill.fore_color.rgb = RGBColor(25, 30, 40)
                insight_bg.line.color.rgb = ACCENT_LIME
                insight_bg.line.width = Pt(1.5)
                
                # Insight sarlavhasi
                insight_header = slide.shapes.add_textbox(Inches(9.3), Inches(4.3), Inches(3.3), Inches(0.4))
                ih_p = insight_header.text_frame.paragraphs[0]
                ih_p.text = "💡 STRATEGIK TAVSIYA"
                ih_p.font.size = Pt(14)
                ih_p.font.bold = True
                ih_p.font.color.rgb = ACCENT_LIME
                
                # Insight matni
                insight_text = slide.shapes.add_textbox(Inches(9.3), Inches(4.7), Inches(3.3), Inches(1.8))
                itf = insight_text.text_frame
                itf.word_wrap = True
                it_p = itf.paragraphs[0]
                it_p.text = insight_val
                it_p.font.size = Pt(15)
                it_p.font.italic = True
                it_p.font.color.rgb = TEXT_WHITE
                it_p.line_spacing = 1.3

            # FOOTER - YANGI DIZAYN
            footer_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 
                0, prs.slide_height - Inches(0.6), prs.slide_width, Inches(0.6))
            footer_bg.fill.solid()
            footer_bg.fill.fore_color.rgb = RGBColor(20, 25, 35)
            footer_bg.line.fill.background()
            
            # Chap tomonda kompaniya nomi
            company_box = slide.shapes.add_textbox(Inches(0.5), prs.slide_height - Inches(0.5), Inches(4), Inches(0.4))
            company_p = company_box.text_frame.paragraphs[0]
            company_p.text = "Slide Master AI"
            company_p.font.size = Pt(12)
            company_p.font.bold = True
            company_p.font.color.rgb = ACCENT_NEON
            
            # O'rtada sanani korsatish
            date_box = slide.shapes.add_textbox(Inches(6), prs.slide_height - Inches(0.5), Inches(2), Inches(0.4))
            date_p = date_box.text_frame.paragraphs[0]
            date_p.text = time.strftime('%Y-%m-%d')
            date_p.font.size = Pt(11)
            date_p.font.color.rgb = TEXT_GRAY
            date_p.alignment = PP_ALIGN.CENTER
            
            # O'ng tomonda slayd raqami
            slide_num_box = slide.shapes.add_textbox(Inches(11.5), prs.slide_height - Inches(0.5), Inches(1.5), Inches(0.4))
            slide_num_p = slide_num_box.text_frame.paragraphs[0]
            slide_num_p.text = f"SLIDE {idx + 1}"
            slide_num_p.font.size = Pt(12)
            slide_num_p.font.bold = True
            slide_num_p.font.color.rgb = ACCENT_CORAL
            slide_num_p.alignment = PP_ALIGN.RIGHT

        os.makedirs("slides", exist_ok=True)
        path = f"slides/Pro_Presentation_{uid}_{int(time.time())}.pptx"
        prs.save(path)
        return path

    except Exception as e:
        logger.error(f"PPTX Generator Error: {e}", exc_info=True)
        return None

# --- 7. HANDLERLAR ---

async def check_sub(user_id):
    try:
        member = await bot.get_chat_member(CHANNEL_ID, user_id)
        return member.status in ['creator', 'administrator', 'member']
    except Exception:
        return True 

async def send_sub_message(message: types.Message, lang):
    btns = [
        [InlineKeyboardButton(text=get_text(lang, 'btn_join'), url=f"https://t.me/{CHANNEL_ID[1:]}")],
        [InlineKeyboardButton(text=get_text(lang, 'btn_check'), callback_data="check_sub")]
    ]
    await message.answer(
        f"{get_text(lang, 'sub_err')}\n\n{CHANNEL_ID}",
        reply_markup=InlineKeyboardMarkup(inline_keyboard=btns),
        parse_mode="Markdown"
    )

async def show_main_menu(message: types.Message, lang):
    b = get_text(lang, 'btns')
    kb = ReplyKeyboardMarkup(
        keyboard=[
            [KeyboardButton(text=b[0]), KeyboardButton(text=b[1])],
            [KeyboardButton(text=b[2]), KeyboardButton(text=b[3])], # Quiz Test 3-indexda
            [KeyboardButton(text=b[4])]
        ],
        resize_keyboard=True
    )
    await message.answer(get_text(lang, 'welcome'), reply_markup=kb)

@dp.message(Command("start"))
async def start_cmd(message: types.Message, command: CommandObject, state: FSMContext):
    await state.clear()
    user = message.from_user
    user_id = user.id
    
    referrer_id = None
    if command.args and command.args.isdigit():
        referrer_id = int(command.args)
        if referrer_id == user_id: referrer_id = None

    is_new = await db.add_user(user_id, user.username, user.first_name, user.last_name, referrer_id)
    
    if is_new and referrer_id:
        await db.update_balance(referrer_id, 1)
        try:
            await bot.send_message(referrer_id, 
                "🎉 **Tabriklaymiz!**\nSizning havolangiz orqali yangi foydalanuvchi qo'shildi.\n💰 Hisobingizga **+1 slayd** qo'shildi!")
        except Exception:
            pass

    user_data = await db.get_user(user_id)
    lang = user_data['lang'] if user_data else 'uz'

    if not await check_sub(user_id):
        return await send_sub_message(message, lang)

    await show_main_menu(message, lang)

@dp.callback_query(F.data == "check_sub")
async def check_sub_callback(callback: CallbackQuery):
    if await check_sub(callback.from_user.id):
        await callback.message.delete()
        user = await db.get_user(callback.from_user.id)
        lang = user['lang'] if user else 'uz'
        await show_main_menu(callback.message, lang)
    else:
        await callback.answer("❌ Hali a'zo bo'lmadingiz!", show_alert=True)

# ----------------- ADMIN COMMANDS -----------------
@dp.message(Command("admin"))
async def admin_panel_cmd(message: types.Message):
    if message.from_user.id != ADMIN_ID: return
    user = await db.get_user(message.from_user.id)
    l = user['lang'] if user else 'uz'
    ikb = InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="📊 Statistika", callback_data="admin_stats")],
        [InlineKeyboardButton(text="📢 Broadcast", callback_data="admin_broadcast")]
    ])
    await message.answer(get_text(l, 'admin_panel'), reply_markup=ikb)

@dp.message(F.text.startswith("/add_"))
async def admin_add_balance_cmd(message: types.Message):
    if message.from_user.id != ADMIN_ID: return
    try:
        parts = message.text.split('_')
        if len(parts) != 3: return await message.answer("❌ Format: /add_USERID_AMOUNT")
        target_id = int(parts[1])
        amount = int(parts[2])
        await db.update_balance(target_id, amount)
        await message.answer(f"✅ User {target_id} balansiga +{amount} slayd qo'shildi!")
    except Exception as e:
        await message.answer(f"❌ Xato: {e}")

@dp.message(F.text.startswith("/vip_"))
async def admin_vip_cmd(message: types.Message):
    if message.from_user.id != ADMIN_ID: return
    try:
        target_id = int(message.text.split('_')[1])
        await db.set_premium(target_id)
        await message.answer(f"✅ User {target_id} VIP Premium statusga o'tkazildi!")
    except Exception as e:
        await message.answer(f"❌ Xato: {e}")

# ----------------- STATE HANDLERS -----------------
@dp.message(UserStates.waiting_package_choice)
async def process_package_choice(message: types.Message, state: FSMContext):
    uid = message.from_user.id
    user = await db.get_user(uid)
    l = user['lang']
    package_btns = get_text(l, 'package_btns')
    text = message.text
    chosen = None
    if text == package_btns[0]: chosen = ("1_slide", 1)
    elif text == package_btns[1]: chosen = ("5_slides", 5)
    elif text == package_btns[2]: chosen = ("vip_premium", 999)
    elif text == get_text(l, 'cancel'):
        await state.clear()
        return await show_main_menu(message, l)
    
    if chosen:
        await state.update_data(chosen_package=chosen[0], amount=chosen[1])
        kb = ReplyKeyboardMarkup(keyboard=[[KeyboardButton(text=get_text(l, 'cancel'))]], resize_keyboard=True)
        await message.answer(get_text(l, 'send_check_now'), reply_markup=kb, parse_mode="Markdown")
        await state.set_state(UserStates.waiting_for_payment)
    else:
        await message.answer(get_text(l, 'choose_package'))

@dp.message(UserStates.waiting_for_payment, F.photo)
async def process_payment(message: types.Message, state: FSMContext):
    if not ADMIN_ID: return await message.answer("❌ Admin sozlanmagan.")
    uid = message.from_user.id
    user = await db.get_user(uid)
    lang = user['lang'] if user else 'uz'
    data = await state.get_data()
    package_type = data.get('chosen_package')
    amount = data.get('amount')

    try:
        payment_id = await db.add_payment(uid, amount, package_type, message.photo[-1].file_id)
        package_names = {'1_slide': "1 ta Slayd (990 so'm)", '5_slides': "5 ta Slayd (2,999 so'm)", 'vip_premium': "VIP Premium (5,999 so'm)"}
        safe_name = html.escape(message.from_user.full_name)
        safe_username = html.escape(message.from_user.username) if message.from_user.username else "Yo'q"
        package_name = package_names.get(package_type, package_type)

        caption = (f"💰 <b>YANGI TO'LOV!</b>\n\n🆔 ID: <code>{uid}</code>\n👤 Foydalanuvchi: {safe_name}\n📱 Username: @{safe_username}\n📦 Paket: <b>{package_name}</b>\n✅ <b>TASDIQLASH UCHUN:</b>\nBalans qo'shish: /add_{uid}_{amount}\nVIP qilish: /vip_{uid}")
        admin_kb = InlineKeyboardMarkup(inline_keyboard=[[InlineKeyboardButton(text="✅ Tasdiqlash", callback_data=f"confirm_{payment_id}"), InlineKeyboardButton(text="❌ Rad etish", callback_data=f"reject_{payment_id}")]])
        await bot.send_photo(chat_id=ADMIN_ID, photo=message.photo[-1].file_id, caption=caption, parse_mode="HTML", reply_markup=admin_kb)
        await message.answer(LANGS[lang]['payment_sent'])
        await state.clear()
        await show_main_menu(message, lang)
    except Exception as e:
        logger.error(f"To'lov yuborishda xato: {e}")
        await message.answer(LANGS[lang]['error'])

@dp.message(AdminStates.waiting_for_broadcast)
async def admin_broadcast_send(message: types.Message, state: FSMContext):
    if message.from_user.id != ADMIN_ID: return
    if message.text and message.text.lower() == "cancel":
        await message.answer("❌ Bekor qilindi.")
        await state.clear()
        return
    
    users = await db.get_all_users()
    count = 0
    await message.answer("⏳ Yuborish boshlandi...")
    for user_row in users:
        try:
            target_id = user_row['id']
            if message.text: await bot.send_message(target_id, message.text, parse_mode="Markdown")
            elif message.photo: await bot.send_photo(target_id, message.photo[-1].file_id, caption=message.caption)
            count += 1
            await asyncio.sleep(0.05)
        except Exception: pass
    await message.answer(f"✅ Xabar {count} ta foydalanuvchiga yuborildi.")
    await state.clear()

# --- QUIZ HANDLER ---
@dp.message(UserStates.waiting_for_quiz_file, F.document)
async def quiz_file_handler(message: types.Message, state: FSMContext):
    uid = message.from_user.id
    user = await db.get_user(uid)
    l = user['lang']

    file_name = message.document.file_name
    file_ext = file_name.split('.')[-1].lower()

    if file_ext not in ['pdf', 'docx', 'txt']:
        return await message.answer("⚠️ Iltimos faqat .PDF, .DOCX yoki .TXT fayl yuboring!")

    await message.answer(get_text(l, 'quiz_processing'))
    await bot.send_chat_action(uid, 'typing')

    file_id = message.document.file_id
    file = await bot.get_file(file_id)
    file_path = file.file_path
    downloaded_file = f"temp_{uid}.{file_ext}"

    try:
        await bot.download_file(file_path, downloaded_file)
        text_content = extract_text_from_file(downloaded_file)

        if not text_content or len(text_content.strip()) < 50:
            os.remove(downloaded_file)
            return await message.answer(get_text(l, 'quiz_error'))

        prompt = (
            f"Analyze the following text and create 10 multiple-choice questions (Quiz). "
            f"Language: {l}. "
            f"Format:\n1. Question?\nA) ...\nB) ...\nC) ...\nAnswer: A\n\n"
            f"Text:\n{text_content}"
        )

        chat_completion = await client.chat.completions.create(
            messages=[
                {"role": "system", "content": "You are a helpful education assistant. Generate a quiz from the provided text."},
                {"role": "user", "content": prompt}
            ],
            model="llama-3.3-70b-versatile",
            temperature=0.5,
            max_tokens=2000
        )

        quiz_result = chat_completion.choices[0].message.content
        
        # Javobni fayl qilib yuborish (uzun bo'lsa)
        if len(quiz_result) > 4000:
            result_file = f"Quiz_{uid}.txt"
            with open(result_file, 'w', encoding='utf-8') as f:
                f.write(quiz_result)
            await bot.send_document(uid, FSInputFile(result_file), caption="✅ **Quiz tayyor!**")
            os.remove(result_file)
        else:
            await message.answer(f"📝 **QUIZ TEST:**\n\n{quiz_result}", parse_mode=None)

    except Exception as e:
        logger.error(f"Quiz Error: {e}")
        await message.answer(get_text(l, 'error'))
    finally:
        if os.path.exists(downloaded_file):
            os.remove(downloaded_file)
        await state.clear()
        await show_main_menu(message, l)


# ----------------- MAIN LOGIC HANDLER -----------------
@dp.message(F.text)
async def main_handler(message: types.Message, state: FSMContext):
    uid = message.from_user.id
    user = await db.get_user(uid)
    if not user: return await message.answer("⚠️ Iltimos /start buyrug'ini bosing.")
    l = user['lang']
    btns = get_text(l, 'btns')
    text = message.text

    if text == btns[0]: # Tariflar
        p_btns = get_text(l, 'package_btns')
        kb = ReplyKeyboardMarkup(keyboard=[[KeyboardButton(text=p_btns[0]), KeyboardButton(text=p_btns[1])],[KeyboardButton(text=p_btns[2])],[KeyboardButton(text=get_text(l, 'cancel'))]], resize_keyboard=True)
        await message.answer(get_text(l, 'tarif'), reply_markup=kb)
        await state.set_state(UserStates.waiting_package_choice)
    elif text == btns[1]: # Kabinet
        status = "⭐ VIP PREMIUM" if user['is_premium'] else "👤 Oddiy"
        msg = (f"📊 **SHAXSIY KABINET**\n\n👤 Ism: {user['first_name']}\n🆔 ID: `{uid}`\n💰 Balans: **{user['balance']} slayd**\n🏷 Status: **{status}**")
        await message.answer(msg, parse_mode="Markdown")
    # main_handler ichida taklif qilish bo'limi
elif text == btns[2]: # Taklif
    bot_info = await bot.get_me()
    link = f"https://t.me/{bot_info.username}?start={uid}"
    
    # YANGI MARKETING MATNI
    promo = f"""🎁 **DO'STLARINGIZNI TAKLIF QILING VA BONUS OLING!**

🔥 Har bir taklif qilingan do'st uchun **+1 BEPUL SLAYD**!
💰 Sizning hisobingizga avtomatik ravishda qo'shiladi.

📲 Do'stlaringizga shu havolani yuboring:

`{link}`

✅ Do'stingiz ro'yxatdan o'tgach, sizning hisobingizga darhol +1 slayd qo'shiladi!

✨ **Bonuslar cheksiz!** Qancha ko'p do'st taklif qilsangiz, shuncha ko'p bepul slaydlar olasiz!"""
    
    kb = ReplyKeyboardMarkup(keyboard=[
        [KeyboardButton(text="📤 Ulashish", request_contact=False)],
        [KeyboardButton(text=get_text(l, 'cancel'))]
    ], resize_keyboard=True)
    
    await message.answer(promo, reply_markup=kb, parse_mode="Markdown")
    elif text == btns[3]: # Quiz Test (YANGI)
        kb = ReplyKeyboardMarkup(keyboard=[[KeyboardButton(text=get_text(l, 'cancel'))]], resize_keyboard=True)
        await message.answer(get_text(l, 'quiz_prompt'), reply_markup=kb, parse_mode="Markdown")
        await state.set_state(UserStates.waiting_for_quiz_file)
    elif text == btns[4]: # Til
        ikb = InlineKeyboardMarkup(inline_keyboard=[[InlineKeyboardButton(text="🇺🇿 O'zbekcha", callback_data="lang_uz")], [InlineKeyboardButton(text="🇷🇺 Русский", callback_data="lang_ru")], [InlineKeyboardButton(text="🇬🇧 English", callback_data="lang_en")]])
        await message.answer("Tilni tanlang / Select language:", reply_markup=ikb)
    elif text == "📤 Ulashish":
    bot_info = await bot.get_me()
    link = f"https://t.me/{bot_info.username}?start={uid}"
    
    # Havolani tezkor nusxa olish uchun tayyor formatda berish
    share_text = f"""🎯 **Slide Master AI - Professional Prezentatsiyalar Boti**

🤖 Sun'iy intellekt yordamida 60 soniyada tayyor prezentatsiyalar!
✅ 3 xil tilda (O'zbek, Rus, Ingliz)
✅ Professional dizayn
✅ PowerPoint formatida

👉 Boshlash uchun: {link}

✨ Do'stlaringizga ulashing va bonuslar oling!"""
    
    await message.answer(share_text, parse_mode="Markdown")
    
    # Havolani alohida ham berish
    await message.answer(f"🔗 **Havola:** `{link}`\n\n📋 *Havolani nusxalash uchun ustiga bosing*", parse_mode="Markdown")
    elif text == get_text(l, 'cancel'):
        await state.clear()
        await show_main_menu(message, l)
    else:
        if not user['is_premium'] and user['balance'] <= 0: return await message.answer(get_text(l, 'no_bal'))
        await state.update_data(topic=text)
        ikb = InlineKeyboardMarkup(inline_keyboard=[[InlineKeyboardButton(text="📄 7 slayd", callback_data="gen:7"), InlineKeyboardButton(text="📄 10 slayd", callback_data="gen:10"), InlineKeyboardButton(text="📄 15 slayd", callback_data="gen:15")]])
        await message.answer(get_text(l, 'gen_prompt').format(topic=text), reply_markup=ikb)

# ----------------- CALLBACK HANDLERS -----------------
@dp.callback_query(F.data.startswith("confirm_"))
async def admin_confirm_payment(callback: CallbackQuery):
    if callback.from_user.id != ADMIN_ID: return
    pid = int(callback.data.split("_")[1])
    async with aiosqlite.connect(DB_PATH) as db_conn:
        db_conn.row_factory = aiosqlite.Row
        cur = await db_conn.execute("SELECT * FROM payments WHERE id = ?", (pid,))
        pay = await cur.fetchone()
        if not pay or pay['status'] != 'pending': return await callback.answer("Eskirgan!", show_alert=True)
        uid, amt, p_type = pay['user_id'], pay['amount'], pay['package_type']
        await db_conn.execute("UPDATE payments SET status = 'approved' WHERE id = ?", (pid,))
        if p_type == 'vip_premium': await db_conn.execute("UPDATE users SET is_premium = 1 WHERE id = ?", (uid,))
        else: await db_conn.execute("UPDATE users SET balance = balance + ? WHERE id = ?", (amt, uid))
        await db_conn.commit()
        try: await bot.send_message(uid, get_text('uz', 'balance_added').format(amount=amt))
        except: pass
    await callback.message.edit_caption(caption=f"✅ Tasdiqlandi!\nID: {pid}")

@dp.callback_query(F.data.startswith("reject_"))
async def admin_reject_payment(callback: CallbackQuery):
    if callback.from_user.id != ADMIN_ID: return
    pid = int(callback.data.split("_")[1])
    async with aiosqlite.connect(DB_PATH) as db_conn:
        await db_conn.execute("UPDATE payments SET status = 'rejected' WHERE id = ?", (pid,))
        await db_conn.commit()
    await callback.message.edit_caption(caption=f"❌ Rad etildi!\nID: {pid}")

@dp.callback_query(F.data.startswith("lang_"))
async def change_lang(callback: CallbackQuery):
    nl = callback.data.split("_")[1]
    await db.update_lang(callback.from_user.id, nl)
    await callback.message.delete()
    await show_main_menu(callback.message, nl)

@dp.callback_query(F.data == "admin_stats")
async def admin_stats_callback(callback: CallbackQuery):
    if callback.from_user.id != ADMIN_ID: return
    st = await db.get_stats()
    await callback.message.answer(f"📊 **STATISTIKA**\n👥 Userlar: {st['total_users']}\n⭐ VIP: {st['premium_users']}\n📈 Slaydlar: {st['total_slides']}", parse_mode="Markdown")

@dp.callback_query(F.data == "admin_broadcast")
async def admin_broadcast_start(callback: CallbackQuery, state: FSMContext):
    if callback.from_user.id != ADMIN_ID: return
    await callback.message.answer("📢 Xabar yuboring:")
    await state.set_state(AdminStates.waiting_for_broadcast)

@dp.callback_query(F.data.startswith("gen:"))
async def generate_ppt(callback: CallbackQuery, state: FSMContext):
    await callback.message.delete()
    uid = callback.from_user.id
    user = await db.get_user(uid)
    l = user['lang']
    if not user['is_premium'] and user['balance'] <= 0: return await callback.message.answer(get_text(l, 'no_bal'))
    
    cnt = callback.data.split(":")[1]
    data = await state.get_data()
    topic = data.get('topic')
    wait_msg = await callback.message.answer(get_text(l, 'wait'))
    await bot.send_chat_action(uid, action="typing")
    path = None
    try:
        lang_instr = {'uz': "IN UZBEK", 'ru': "IN RUSSIAN", 'en': "IN ENGLISH"}.get(l, "IN UZBEK")
        sys_p = ("You are a Senior Presentation Consultant. "
                 f"STRICTLY {lang_instr}. Create dense content. "
                 "Return ONLY VALID JSON: {'slides': [{'title': '...', 'content': [{'bold': '...', 'text': '...'}], 'stat': '...', 'insight': '...'}]}")
        usr_p = f"Create {cnt}-slide presentation on '{topic}'."
        
        comp = await client.chat.completions.create(messages=[{"role":"system","content":sys_p},{"role":"user","content":usr_p}], model="llama-3.3-70b-versatile", response_format={"type":"json_object"})
        path = await asyncio.to_thread(create_ultra_modern_pptx, topic, comp.choices[0].message.content, uid)
        
        if path:
            await bot.send_document(uid, FSInputFile(path), caption=get_text(l, 'done'))
            if not user['is_premium']: await db.update_balance(uid, -1)
        else: await callback.message.answer(get_text(l, 'error'))
    except Exception as e:
        logger.error(f"Gen Error: {e}")
        await callback.message.answer(get_text(l, 'error'))
    finally:
        try: await wait_msg.delete()
        except: pass
        await state.clear()
        if path and os.path.exists(path):
            await asyncio.sleep(2)
            try: os.remove(path)
            except: pass

async def main():
    await db.init()
    os.makedirs("slides", exist_ok=True)
    try: await bot.delete_webhook(drop_pending_updates=True)
    except: pass
    await dp.start_polling(bot)

if __name__ == "__main__":
    try: asyncio.run(main())
    except KeyboardInterrupt: pass
